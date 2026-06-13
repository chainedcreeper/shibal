"""PPTX 텍스트 추출 — 라이브러리 병렬.

- python-pptx (공식 API: shapes/text_frame 순회)
- zipfile + lxml 로 ppt/slides/slide*.xml 직접 파싱 (폴백)
"""
from __future__ import annotations

import zipfile

from ._parallel import run_parallel


def _python_pptx(path: str) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(path)
    pages: list[dict] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(r.text for r in para.runs).strip()
                if line:
                    texts.append(line)
            # 노트도 있으면 추가
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[노트] {notes}")
        if texts:
            pages.append({"text": "\n".join(texts), "page": i})
    return pages


def _xml_raw_pptx(path: str) -> list[dict]:
    from lxml import etree

    pages: list[dict] = []
    with zipfile.ZipFile(path) as zf:
        names = sorted(
            n for n in zf.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
        for i, name in enumerate(names, 1):
            with zf.open(name) as f:
                tree = etree.parse(f)
            texts = [
                (t.text or "").strip()
                for t in tree.iter()
                if isinstance(t.tag, str) and t.tag.rsplit("}", 1)[-1] in ("t", "a:t")
            ]
            joined = "\n".join(t for t in texts if t)
            if joined:
                pages.append({"text": joined, "page": i})
    return pages


_EXTRACTORS = {
    "python-pptx": _python_pptx,
    "xml-raw":     _xml_raw_pptx,
}


def extract_text(path: str) -> list[dict]:
    return run_parallel(_EXTRACTORS, path, min_chars=20)
