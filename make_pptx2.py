"""
머신러닝 해보기 세미나 PPT — 9장
python make_pptx2.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG_DARK  = RGBColor(0x0f, 0x0f, 0x1a)
BG_SLIDE = RGBColor(0x11, 0x11, 0x22)
BLUE     = RGBColor(0x7e, 0xb8, 0xff)
BLUE_MID = RGBColor(0x3a, 0x7a, 0xff)
TEXT     = RGBColor(0xe8, 0xe8, 0xf0)
TEXT_SUB = RGBColor(0xc8, 0xd8, 0xf0)
DIM      = RGBColor(0x70, 0x80, 0xa0)
GREEN    = RGBColor(0x5d, 0xdf, 0x5d)
ORANGE   = RGBColor(0xff, 0xaa, 0x40)
PURPLE   = RGBColor(0xd2, 0xa8, 0xff)
TEAL     = RGBColor(0x4a, 0xd4, 0xc0)
YELLOW   = RGBColor(0xff, 0xd0, 0x50)
RED      = RGBColor(0xff, 0x70, 0x70)
CARD     = RGBColor(0x1a, 0x1e, 0x30)
BORDER   = RGBColor(0x2a, 0x35, 0x60)

W = Inches(13.33)
H = Inches(7.5)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for ph in slide.placeholders:
        ph._element.getparent().remove(ph._element)
    return slide

def bg(slide, color=BG_SLIDE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def T(slide, t, l, top, w, h,
      size=Pt(16), bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, top, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = t; r.font.size = size
    r.font.bold = bold; r.font.color.rgb = color
    return tb

def ML(slide, items, l, top, w, h, size=Pt(14), color=TEXT_SUB):
    tb = slide.shapes.add_textbox(l, top, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            t, b, c = item, False, color
        else:
            t = item[0]
            b = item[1] if len(item) > 1 else False
            c = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_after = Pt(6)
        r = p.add_run()
        r.text = t; r.font.size = size
        r.font.bold = b; r.font.color.rgb = c

def hdr(slide, title, subtitle=None):
    bg(slide)
    box(slide, 0, 0, W, Inches(0.1), fill=BLUE_MID)
    T(slide, title,
      Inches(0.7), Inches(0.2), Inches(11.9), Inches(0.52),
      size=Pt(28), bold=True, color=TEXT)
    if subtitle:
        T(slide, subtitle,
          Inches(0.7), Inches(0.74), Inches(11.9), Inches(0.3),
          size=Pt(14), color=DIM)
        box(slide, Inches(0.7), Inches(1.08), Inches(12.0), Pt(1.2), fill=BLUE_MID)
        return Inches(1.25)
    box(slide, Inches(0.7), Inches(0.76), Inches(12.0), Pt(1.2), fill=BLUE_MID)
    return Inches(0.94)


# ── S1 타이틀 ─────────────────────────────────────────────────────────
def s1(prs):
    slide = new_slide(prs)
    bg(slide, BG_DARK)
    box(slide, 0, 0, W, Inches(0.18), fill=BLUE_MID)
    box(slide, 0, Inches(7.32), W, Inches(0.18), fill=BLUE_MID)

    T(slide, "머신러닝 해보기",
      Inches(1.4), Inches(2.0), Inches(10.5), Inches(1.2),
      size=Pt(64), bold=True, color=BLUE)
    T(slide, "LLM 파인튜닝 시도기  —  Knowledge Distillation + LoRA",
      Inches(1.4), Inches(3.3), Inches(10.5), Inches(0.45),
      size=Pt(18), color=DIM)
    T(slide, "2026. 06. 12",
      Inches(1.4), Inches(6.88), Inches(4), Inches(0.35),
      size=Pt(13), color=DIM)


# ── S2 파인튜닝이란? ─────────────────────────────────────────────────
def s2(prs):
    slide = new_slide(prs)
    y0 = hdr(slide, "파인튜닝이란?", "Fine-tuning")

    # 정의
    box(slide, Inches(0.7), y0+Inches(0.18),
        Inches(12.0), Inches(1.1),
        fill=RGBColor(0x0f,0x1a,0x30), line=BLUE_MID)
    T(slide,
      "사전학습된 LLM을 우리 목적에 맞게 추가 학습하는 것",
      Inches(0.95), y0+Inches(0.35),
      Inches(11.5), Inches(0.55),
      size=Pt(22), bold=True, color=BLUE)

    # 왼쪽: 왜 필요한가
    box(slide, Inches(0.7), y0+Inches(1.45),
        Inches(5.7), Inches(4.3),
        fill=RGBColor(0x25,0x10,0x10), line=RED, lw=Pt(1.5))
    T(slide, "범용 LLM의 한계",
      Inches(0.95), y0+Inches(1.62),
      Inches(5.2), Inches(0.4),
      size=Pt(18), bold=True, color=RED)
    ML(slide, [
        "•  도메인 지식 부족",
        "•  할루시네이션  (지문 벗어난 내용 지어냄)",
        "•  추론 없이 단답만 뱉음",
        "•  서버 24시간 운용 현실적으로 불가",
    ], Inches(0.95), y0+Inches(2.1),
       Inches(5.2), Inches(3.0),
       size=Pt(16), color=TEXT_SUB)

    # 오른쪽: 우리가 선택한 방식
    box(slide, Inches(6.73), y0+Inches(1.45),
        Inches(5.97), Inches(4.3),
        fill=RGBColor(0x11,0x1e,0x11), line=GREEN, lw=Pt(1.5))
    T(slide, "우리가 선택한 방식",
      Inches(6.98), y0+Inches(1.62),
      Inches(5.47), Inches(0.4),
      size=Pt(18), bold=True, color=GREEN)
    ML(slide, [
        ("Knowledge Distillation", True, ORANGE),
        "  큰 모델(Teacher)의 추론 과정을",
        "  작은 모델(Student)에 이식",
        "",
        ("LoRA", True, GREEN),
        "  모델 전체가 아닌 작은 보정 행렬만 학습",
        "  48GB VRAM으로 충분히 가능",
    ], Inches(6.98), y0+Inches(2.1),
       Inches(5.47), Inches(3.0),
       size=Pt(16), color=TEXT_SUB)


# ── S4 파인튜닝이 뭔지 ────────────────────────────────────────────────
def s4(prs):
    slide = new_slide(prs)
    y0 = hdr(slide, "파인튜닝 — Knowledge Distillation + LoRA")

    # KD
    box(slide, Inches(0.7), y0+Inches(0.15),
        Inches(5.87), Inches(5.5),
        fill=RGBColor(0x1a,0x14,0x0a), line=RGBColor(0x5a,0x3a,0x10))
    T(slide, "Knowledge Distillation",
      Inches(0.95), y0+Inches(0.32),
      Inches(5.37), Inches(0.48),
      size=Pt(22), bold=True, color=ORANGE)
    T(slide, "지식 증류",
      Inches(0.95), y0+Inches(0.82),
      Inches(5.37), Inches(0.3),
      size=Pt(14), color=DIM)
    box(slide, Inches(0.95), y0+Inches(1.18),
        Inches(5.37), Pt(1), fill=ORANGE)
    ML(slide, [
        "큰 모델(Teacher)의 지식을",
        "작은 모델(Student)에 이식",
        "",
        ("Teacher  :  Qwen3-14B", False, ORANGE),
        ("Student  :  Qwen3-8B", False, TEXT_SUB),
        "",
        "단순 정답이 아니라",
        "<think> 추론 과정까지 학습",
        "→  Reasoning Distillation",
        "",
        "DeepSeek-R1, Qwen3 distill도",
        "이 방식으로 만들어진 모델",
    ], Inches(0.95), y0+Inches(1.35),
       Inches(5.37), Inches(3.8),
       size=Pt(15), color=TEXT_SUB)

    # LoRA
    box(slide, Inches(6.83), y0+Inches(0.15),
        Inches(5.87), Inches(5.5),
        fill=RGBColor(0x11,0x1e,0x11), line=RGBColor(0x2a,0x5c,0x2a))
    T(slide, "LoRA",
      Inches(7.08), y0+Inches(0.32),
      Inches(5.37), Inches(0.48),
      size=Pt(22), bold=True, color=GREEN)
    T(slide, "Low-Rank Adaptation",
      Inches(7.08), y0+Inches(0.82),
      Inches(5.37), Inches(0.3),
      size=Pt(14), color=DIM)
    box(slide, Inches(7.08), y0+Inches(1.18),
        Inches(5.37), Pt(1), fill=GREEN)
    ML(slide, [
        "모델 전체를 학습하는 대신",
        "내부 행렬 옆에 작은 보정 행렬",
        "두 개(A, B)만 추가해서 학습",
        "",
        ("W'  =  W₀  +  B · A", True, YELLOW),
        "",
        "W₀는 고정  →  A, B만 업데이트",
        "전체 파라미터의 0.3%만 학습",
        "",
        ("Full FT  :  ~112GB VRAM", False, RED),
        ("LoRA     :  ~28GB VRAM", False, GREEN),
        "",
        "48GB 서버에서 가능한 이유",
    ], Inches(7.08), y0+Inches(1.35),
       Inches(5.37), Inches(3.8),
       size=Pt(15), color=TEXT_SUB)


# ── S5 어떻게 할 건지 ─────────────────────────────────────────────────
def s5(prs):
    slide = new_slide(prs)
    y0 = hdr(slide, "어떻게 할 건지 — 전체 파이프라인")

    steps = [
        ("Wikipedia\n한국어", "지문 소스",    BLUE,   RGBColor(0x1a,0x1a,0x3a)),
        ("Teacher\nQwen3-14B", "CoT 생성",   ORANGE, RGBColor(0x2a,0x18,0x08)),
        ("데이터셋\n병합",     "~70,000개",   GREEN,  RGBColor(0x11,0x1e,0x11)),
        ("Student\nQwen3-8B", "LoRA 학습",   PURPLE, RGBColor(0x1a,0x11,0x2a)),
        ("GGUF\n변환",        "양자화",       TEAL,   RGBColor(0x0f,0x1e,0x1e)),
        ("Ollama\n서빙",      "로컬 배포",    YELLOW, RGBColor(0x1e,0x1e,0x08)),
    ]
    sw = Inches(1.87); sh = Inches(1.3)
    sx = Inches(0.7)
    for i, (label, sub, acc, bg_c) in enumerate(steps):
        box(slide, sx, y0+Inches(0.18), sw, sh, fill=bg_c, line=acc, lw=Pt(2))
        T(slide, label, sx, y0+Inches(0.3), sw, Inches(0.65),
          size=Pt(14), bold=True, color=acc, align=PP_ALIGN.CENTER)
        T(slide, sub, sx, y0+Inches(0.98), sw, Inches(0.35),
          size=Pt(11), color=DIM, align=PP_ALIGN.CENTER)
        sx += sw
        if i < len(steps)-1:
            T(slide, "→", sx, y0+Inches(0.42), Inches(0.28), Inches(0.45),
              size=Pt(20), color=DIM, align=PP_ALIGN.CENTER)
            sx += Inches(0.28)

    y1 = y0+Inches(1.68)
    cols = [
        ("데이터", ORANGE,
         ["Wikipedia CoT   5,000개",
          "KorQuAD 1.0    10,000개",
          "KLUE / MRC     10,000개",
          "sae4K          10,000개",
          "OpenThoughts   30,000개",
          "",
          ("합계  ~65,000개", True, ORANGE),
          "",
          "소요 시간",
          ("~13시간 (현재 진행 중)", False, DIM),
         ]),
        ("학습", GREEN,
         ["Qwen3-8B  LoRA r=32",
          "fp16  ·  에폭 2",
          "Kubeflow L40S 48GB",
          "",
          "Teacher <think> 추론 과정",
          "포함 학습",
          "",
          "소요 시간",
          ("~2~3시간", True, GREEN),
         ]),
        ("배포", TEAL,
         ["GGUF 변환  (q4_k_m)",
          "Ollama 등록",
          "Tailscale VPN 연결",
          "AI Tutor 모델명 교체",
          "",
          "소요 시간",
          ("~30분", True, TEAL),
         ]),
    ]
    cw = Inches(3.93); ch = Inches(4.42)
    cx = Inches(0.7)
    for title, acc, bullets in cols:
        box(slide, cx, y1, cw, ch, fill=CARD, line=acc, lw=Pt(1.5))
        T(slide, title, cx+Inches(0.2), y1+Inches(0.15),
          cw-Inches(0.4), Inches(0.42),
          size=Pt(17), bold=True, color=acc)
        box(slide, cx+Inches(0.2), y1+Inches(0.62),
            cw-Inches(0.4), Pt(1), fill=acc)
        ML(slide, bullets,
           cx+Inches(0.2), y1+Inches(0.77),
           cw-Inches(0.4), ch-Inches(0.95),
           size=Pt(14), color=TEXT_SUB)
        cx += cw + Inches(0.17)


# ── S6 지금 하고 있는 것 ──────────────────────────────────────────────
def s6(prs):
    slide = new_slide(prs)
    y0 = hdr(slide, "지금 하고 있는 것")

    box(slide, Inches(0.7), y0+Inches(0.2),
        Inches(12.0), Inches(5.5),
        fill=CARD, line=BORDER)

    # 진행 상태
    status = [
        ("✅", "Wikipedia CoT 데이터 생성 중",
         "Kubeflow L40S 48GB  ·  목표 5,000개  ·  현재 ~50% 완료",
         "소요: ~13시간  (진행 중)", GREEN),
        ("⏳", "데이터셋 병합",
         "KorQuAD + KLUE/MRC + OpenThoughts + sae4K  →  ~65,000개",
         "소요: ~30분", ORANGE),
        ("□",  "Student 모델 학습",
         "Qwen3-8B  LoRA r=32  ·  fp16  ·  에폭 2",
         "소요: ~2~3시간", DIM),
        ("□",  "평가  (LLM-as-a-Judge)",
         "Teacher가 Student 답변 검증  ·  샘플 200개 기준",
         "소요: ~3~4시간  (결과 따라 재학습 시 +2~3시간)", DIM),
        ("□",  "GGUF 변환 + 배포",
         "llama.cpp 변환  →  Ollama 등록  →  AI Tutor 연결",
         "소요: ~30분", DIM),
    ]
    for i, (icon, title, desc, time_str, acc) in enumerate(status):
        ry = y0+Inches(0.4)+Inches(1.02)*i
        T(slide, icon,
          Inches(0.95), ry+Inches(0.1), Inches(0.5), Inches(0.55),
          size=Pt(20), color=acc, align=PP_ALIGN.CENTER)
        T(slide, title,
          Inches(1.6), ry+Inches(0.05),
          Inches(8.5), Inches(0.35),
          size=Pt(16), bold=True, color=acc)
        T(slide, desc,
          Inches(1.6), ry+Inches(0.4),
          Inches(8.5), Inches(0.3),
          size=Pt(12), color=DIM)
        T(slide, time_str,
          Inches(10.2), ry+Inches(0.15),
          Inches(2.8), Inches(0.35),
          size=Pt(12), bold=True, color=acc, align=PP_ALIGN.RIGHT)


# ── S7 추가 데이터셋 ──────────────────────────────────────────────────
def s7(prs):
    slide = new_slide(prs)
    y0 = hdr(slide, "추가로 고려한 데이터셋", "AwesomeKorean_Data 등에서 찾음")

    datasets = [
        ("Wikipedia CoT", "직접 생성", ORANGE, RGBColor(0x2a,0x18,0x08), RGBColor(0x5a,0x3a,0x10),
         ["Teacher가 직접 생성",
          "도메인 맞춤 CoT 포함",
          "목표 5,000개",
          "",
          "핵심 데이터"]),
        ("KorQuAD 1.0", "squad_kor_v1", BLUE, RGBColor(0x0f,0x1a,0x30), BORDER,
         ["한국어 독해 QA",
          "약 60,000개",
          "위키피디아 기반",
          "",
          "지문 충실도 강화"]),
        ("KLUE / MRC", "klue/mrc", TEAL, RGBColor(0x0f,0x1e,0x1e), RGBColor(0x2a,0x5c,0x5c),
         ["한국어 독해 QA",
          "약 10,000개",
          "다양한 도메인",
          "",
          "한국어 이해력 향상"]),
        ("sae4K", "naver-news-summarization", GREEN, RGBColor(0x11,0x1e,0x11), RGBColor(0x2a,0x5c,0x2a),
         ["한국어 뉴스 요약",
          "약 50,000개",
          "문서→요약 쌍",
          "",
          "요약 능력 학습"]),
        ("OpenThoughts\n-114k", "open-thoughts", PURPLE, RGBColor(0x1a,0x11,0x2a), RGBColor(0x5a,0x2a,0x8a),
         ["영어 CoT 추론",
          "30,000개 사용",
          "수학·논리 추론",
          "",
          "추론 능력 보강"]),
    ]

    cw = Inches(2.3); ch = Inches(5.3)
    cx = Inches(0.7)
    for name, repo, acc, bg_c, bd, bullets in datasets:
        box(slide, cx, y0+Inches(0.15), cw, ch, fill=bg_c, line=bd, lw=Pt(1.5))
        T(slide, name,
          cx+Inches(0.15), y0+Inches(0.28),
          cw-Inches(0.3), Inches(0.65),
          size=Pt(15), bold=True, color=acc, align=PP_ALIGN.CENTER)
        T(slide, repo,
          cx+Inches(0.15), y0+Inches(0.95),
          cw-Inches(0.3), Inches(0.28),
          size=Pt(10), color=DIM, align=PP_ALIGN.CENTER)
        box(slide, cx+Inches(0.15), y0+Inches(1.28),
            cw-Inches(0.3), Pt(1), fill=acc)
        ML(slide, bullets,
           cx+Inches(0.15), y0+Inches(1.45),
           cw-Inches(0.3), ch-Inches(1.6),
           size=Pt(13), color=TEXT_SUB)
        cx += cw + Inches(0.2)

    by = y0+Inches(0.15)+ch+Inches(0.2)
    box(slide, Inches(0.7), by, Inches(12.0), Inches(0.55),
        fill=RGBColor(0x0f,0x1a,0x30), line=BLUE_MID)
    T(slide, "합계  :  Wikipedia 5k  +  KorQuAD 10k  +  KLUE 10k  +  sae4K 10k  +  OpenThoughts 30k  =  ~65,000개",
      Inches(0.95), by+Inches(0.1), Inches(11.5), Inches(0.38),
      size=Pt(14), bold=True, color=BLUE)


# ── S8 평가 방법 ──────────────────────────────────────────────────────
def s8(prs):
    slide = new_slide(prs)
    y0 = hdr(slide, "평가를 어떻게 할 것인가", "고민 중")

    # 메인: LLM-as-a-Judge
    box(slide, Inches(0.7), y0+Inches(0.18),
        Inches(7.5), Inches(5.5),
        fill=RGBColor(0x0f,0x1a,0x30), line=BLUE_MID, lw=Pt(2))
    T(slide, "LLM-as-a-Judge  (메인)",
      Inches(0.95), y0+Inches(0.35),
      Inches(7.0), Inches(0.45),
      size=Pt(20), bold=True, color=BLUE)
    box(slide, Inches(0.95), y0+Inches(0.85),
        Inches(7.0), Pt(1), fill=BLUE_MID)
    ML(slide, [
        "Teacher(Qwen3-14B)가 채점관이 되어",
        "Student(Qwen3-8B)의 답변을 직접 평가",
        "",
        ("평가 기준", True, BLUE),
        "  •  지문 근거했는가  (Faithfulness)",
        "  •  추론 과정이 올바른가  (CoT 품질)",
        "  •  Teacher 답변과 얼마나 유사한가",
        "",
        ("선택 이유", True, BLUE),
        "  •  도메인에 맞는 평가 가능",
        "  •  숫자 지표보다 직관적",
        "  •  DeepSeek, Qwen 계열 평가에서",
        "     실제로 많이 쓰이는 방식",
    ], Inches(0.95), y0+Inches(1.02),
       Inches(7.0), Inches(4.3),
       size=Pt(15), color=TEXT_SUB)

    # 보조 지표
    box(slide, Inches(8.43), y0+Inches(0.18),
        Inches(4.3), Inches(5.5),
        fill=CARD, line=BORDER)
    T(slide, "보조 지표",
      Inches(8.68), y0+Inches(0.35),
      Inches(3.8), Inches(0.42),
      size=Pt(18), bold=True, color=DIM)
    box(slide, Inches(8.68), y0+Inches(0.82),
        Inches(3.8), Pt(1), fill=BORDER)
    metrics = [
        ("ROUGE-L",      "텍스트 겹침 F1",    "≥ 0.35"),
        ("BERTScore",    "의미적 유사도",      "≥ 0.80"),
        ("Faithfulness", "지문 충실도",        "≥ 0.70"),
        ("Hallucination","사실 날조 비율",     "≤ 0.10"),
    ]
    for i, (name, desc, target) in enumerate(metrics):
        ry = y0+Inches(1.05)+Inches(1.08)*i
        T(slide, name, Inches(8.68), ry,
          Inches(2.0), Inches(0.38),
          size=Pt(14), bold=True, color=TEXT)
        T(slide, desc, Inches(8.68), ry+Inches(0.38),
          Inches(2.5), Inches(0.3),
          size=Pt(12), color=DIM)
        T(slide, target, Inches(10.8), ry+Inches(0.05),
          Inches(1.7), Inches(0.38),
          size=Pt(14), bold=True, color=GREEN)


# ── S9 마무리 ─────────────────────────────────────────────────────────
def s9(prs):
    slide = new_slide(prs)
    bg(slide, BG_DARK)
    box(slide, 0, 0, W, Inches(0.18), fill=BLUE_MID)
    box(slide, 0, Inches(7.32), W, Inches(0.18), fill=BLUE_MID)

    T(slide, "지금까지",
      Inches(1.4), Inches(0.45),
      Inches(10), Inches(0.5),
      size=Pt(22), bold=True, color=BLUE)

    points = [
        ("계기",     "48GB VRAM 서버  +  AI Tutor 할루시네이션 문제",     BLUE),
        ("선택",     "Knowledge Distillation  +  LoRA 파인튜닝",          ORANGE),
        ("진행",     "Teacher(14B) → CoT 데이터 생성 중  (현재 ~50%)",    GREEN),
        ("데이터",   "Wikipedia + KorQuAD + KLUE + sae4K  ~65,000개",     TEAL),
        ("평가",     "LLM-as-a-Judge  —  Teacher가 Student 검증",         PURPLE),
        ("목표",     "작은 모델로 로컬에서 큰 모델 수준 내기",             YELLOW),
    ]

    for i, (label, desc, acc) in enumerate(points):
        lx = Inches(1.4) if i < 3 else Inches(7.2)
        ly = Inches(1.15) + Inches(0.95)*(i % 3)
        box(slide, lx, ly, Inches(5.5), Inches(0.75),
            fill=CARD, line=acc, lw=Pt(1.5))
        T(slide, label,
          lx+Inches(0.2), ly+Inches(0.1),
          Inches(1.2), Inches(0.55),
          size=Pt(15), bold=True, color=acc)
        T(slide, desc,
          lx+Inches(1.5), ly+Inches(0.15),
          Inches(3.8), Inches(0.48),
          size=Pt(13.5), color=TEXT_SUB)

    T(slide, "Q & A",
      Inches(1.4), Inches(4.55),
      Inches(10.5), Inches(0.9),
      size=Pt(52), bold=True, color=BLUE,
      align=PP_ALIGN.CENTER)


# ── 실행 ──────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("생성 중...")
    s1(prs); print("  1/7 타이틀")
    s2(prs); print("  2/7 파인튜닝이란?")
    s4(prs); print("  3/7 KD + LoRA 상세")
    s5(prs); print("  4/7 어떻게 할 건지 + 데이터셋")
    s6(prs); print("  5/7 지금 하고 있는 것")
    s8(prs); print("  6/7 평가 방법")
    s9(prs); print("  7/7 마무리")

    prs.save("finetune_seminar_v2.pptx")
    print("\n완료: finetune_seminar_v2.pptx")

if __name__ == "__main__":
    main()
